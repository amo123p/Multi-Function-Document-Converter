#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
多功能文档转换工具 v4.2
- 支持批量多文件选择
- 修复浏览器驱动问题，支持离线使用
- 完整保留原始样式
"""

import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from tkinter.scrolledtext import ScrolledText
import os
import sys
import threading
import tempfile
import subprocess
import time
import shutil
import winreg
from pathlib import Path
from io import BytesIO

# ============== 检查依赖 ==============
def check_dependencies():
    required = {
        'PIL': 'Pillow',
        'pptx': 'python-pptx',
        'fitz': 'PyMuPDF',
    }
    missing = []
    for module, package in required.items():
        try:
            __import__(module)
        except ImportError:
            missing.append(package)
    
    if missing:
        print(f"❌ 缺少依赖包: {', '.join(missing)}")
        print(f"请运行: pip install {' '.join(missing)}")
        return False
    return True

if not check_dependencies():
    input("按回车键退出...")
    sys.exit(1)

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu
import fitz

# 可选依赖
try:
    import win32com.client
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False
    print("⚠️ 未安装 pywin32，请运行: pip install pywin32")

try:
    from selenium import webdriver
    from selenium.webdriver.chrome.service import Service as ChromeService
    from selenium.webdriver.edge.service import Service as EdgeService
    from selenium.webdriver.chrome.options import Options as ChromeOptions
    from selenium.webdriver.edge.options import Options as EdgeOptions
    from selenium.webdriver.common.by import By
    from selenium.webdriver.support.ui import WebDriverWait
    from selenium.webdriver.support import expected_conditions as EC
    HAS_SELENIUM = True
except ImportError:
    HAS_SELENIUM = False
    print("⚠️ 未安装 selenium，请运行: pip install selenium")


class TaskController:
    """任务控制器"""
    def __init__(self):
        self.pause_event = threading.Event()
        self.pause_event.set()
        self.stop_flag = False
        self.is_running = False
    
    def pause(self):
        self.pause_event.clear()
    
    def resume(self):
        self.pause_event.set()
    
    def stop(self):
        self.stop_flag = True
        self.pause_event.set()
    
    def reset(self):
        self.pause_event.set()
        self.stop_flag = False
        self.is_running = False
    
    def check_pause(self):
        self.pause_event.wait()
        return not self.stop_flag
    
    def should_stop(self):
        return self.stop_flag


class BrowserDriverManager:
    """浏览器驱动管理器 - 支持离线使用"""
    
    def __init__(self, log_callback=None):
        self.log = log_callback or print
        self.driver_paths = {
            'edge': None,
            'chrome': None
        }
        self._find_drivers()
    
    def _find_drivers(self):
        """查找本地驱动"""
        search_paths = [
            os.getcwd(),
            os.path.dirname(os.path.abspath(__file__)),
            r"C:\WebDriver",
            r"C:\Drivers",
            os.path.expanduser("~"),
            os.path.expanduser("~/Downloads"),
            os.path.expanduser("~/Desktop"),
            r"C:\Program Files\WebDriver",
            r"C:\Program Files (x86)\WebDriver",
        ]
        
        path_env = os.environ.get('PATH', '')
        search_paths.extend(path_env.split(os.pathsep))
        
        edge_names = ['msedgedriver.exe', 'MicrosoftWebDriver.exe']
        for path in search_paths:
            if not os.path.exists(path):
                continue
            for name in edge_names:
                driver_path = os.path.join(path, name)
                if os.path.exists(driver_path):
                    self.driver_paths['edge'] = driver_path
                    self.log(f"✅ 找到Edge驱动: {driver_path}")
                    break
            if self.driver_paths['edge']:
                break
        
        chrome_names = ['chromedriver.exe', 'chromedriver']
        for path in search_paths:
            if not os.path.exists(path):
                continue
            for name in chrome_names:
                driver_path = os.path.join(path, name)
                if os.path.exists(driver_path):
                    self.driver_paths['chrome'] = driver_path
                    self.log(f"✅ 找到Chrome驱动: {driver_path}")
                    break
            if self.driver_paths['chrome']:
                break
    
    def get_edge_driver(self):
        return self.driver_paths.get('edge')
    
    def get_chrome_driver(self):
        return self.driver_paths.get('chrome')
    
    def get_browser_version(self, browser='edge'):
        try:
            if browser == 'edge':
                key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, 
                    r"Software\Microsoft\Edge\BLBeacon")
                version, _ = winreg.QueryValueEx(key, "version")
                return version
            elif browser == 'chrome':
                key = winreg.OpenKey(winreg.HKEY_CURRENT_USER,
                    r"Software\Google\Chrome\BLBeacon")
                version, _ = winreg.QueryValueEx(key, "version")
                return version
        except:
            return None
    
    def download_driver_instructions(self):
        edge_ver = self.get_browser_version('edge') or "未知"
        chrome_ver = self.get_browser_version('chrome') or "未知"
        
        return f"""
════════════════════════════════════════════════════════
  浏览器驱动下载说明
════════════════════════════════════════════════════════

您的浏览器版本:
  • Edge: {edge_ver}
  • Chrome: {chrome_ver}

下载地址:
  • Edge驱动: https://developer.microsoft.com/en-us/microsoft-edge/tools/webdriver/
  • Chrome驱动: https://googlechromelabs.github.io/chrome-for-testing/

安装步骤:
  1. 根据您的浏览器版本下载对应的驱动
  2. 解压下载的文件
  3. 将 msedgedriver.exe 或 chromedriver.exe 放到以下位置之一:
     - 程序所在目录: {os.getcwd()}
     - C:\\WebDriver\\
     - 系统PATH目录

════════════════════════════════════════════════════════
"""


class DocumentConverter:
    """文档转换核心类"""
    
    def __init__(self, log_callback=None, progress_callback=None):
        self.log = log_callback or print
        self.progress = progress_callback or (lambda x, y: None)
        self.controller = TaskController()
        self.driver_manager = BrowserDriverManager(self.log)
        self._check_available_tools()
    
    def _check_available_tools(self):
        """检查可用工具"""
        self.tools = {
            'ms_word': False,
            'ms_excel': False,
            'wps': False,
            'libreoffice': False,
            'chrome': None,
            'edge': None,
        }
        
        if HAS_WIN32COM:
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Quit()
                self.tools['ms_word'] = True
                self.log("✅ 检测到 Microsoft Word")
            except:
                pass
            
            try:
                excel = win32com.client.Dispatch("Excel.Application")
                excel.Quit()
                self.tools['ms_excel'] = True
                self.log("✅ 检测到 Microsoft Excel")
            except:
                pass
            
            try:
                wps = win32com.client.Dispatch("KWPS.Application")
                wps.Quit()
                self.tools['wps'] = True
                self.log("✅ 检测到 WPS Office")
            except:
                try:
                    wps = win32com.client.Dispatch("KET.Application")
                    wps.Quit()
                    self.tools['wps'] = True
                except:
                    pass
        
        libreoffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for path in libreoffice_paths:
            if os.path.exists(path):
                self.tools['libreoffice'] = path
                self.log("✅ 检测到 LibreOffice")
                break
        
        edge_paths = [
            r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
            r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
        ]
        for path in edge_paths:
            if os.path.exists(path):
                self.tools['edge'] = path
                self.log("✅ 检测到 Edge 浏览器")
                break
        
        chrome_paths = [
            r"C:\Program Files\Google\Chrome\Application\chrome.exe",
            r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
            os.path.expandvars(r"%LOCALAPPDATA%\Google\Chrome\Application\chrome.exe"),
        ]
        for path in chrome_paths:
            if os.path.exists(path):
                self.tools['chrome'] = path
                self.log("✅ 检测到 Chrome 浏览器")
                break
        
        if self.driver_manager.get_edge_driver():
            self.log("✅ Edge驱动已就绪")
        elif self.tools['edge']:
            self.log("⚠️ 未找到Edge驱动，网页转换需要下载驱动")
        
        if self.driver_manager.get_chrome_driver():
            self.log("✅ Chrome驱动已就绪")
        elif self.tools['chrome']:
            self.log("⚠️ 未找到Chrome驱动")
    
    def get_controller(self):
        return self.controller
    
    def get_driver_instructions(self):
        return self.driver_manager.download_driver_instructions()
    
    # ==================== 批量文档转PDF ====================
    def documents_to_pdf(self, doc_paths, output_folder):
        """批量文档转PDF"""
        try:
            self.controller.is_running = True
            total = len(doc_paths)
            success_count = 0
            
            self.log(f"🔄 批量转换 {total} 个文档...")
            os.makedirs(output_folder, exist_ok=True)
            
            for i, doc_path in enumerate(doc_paths):
                if not self.controller.check_pause():
                    break
                
                self.progress(i + 1, total)
                base_name = Path(doc_path).stem
                output_path = os.path.join(output_folder, f"{base_name}.pdf")
                
                if self.document_to_pdf(doc_path, output_path):
                    success_count += 1
            
            self.log(f"✅ 完成！成功 {success_count}/{total}")
            return success_count > 0
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
        finally:
            self.controller.is_running = False
    
    def document_to_pdf(self, doc_path, output_path):
        """单个文档转PDF"""
        try:
            self.log(f"🔄 转换文档: {os.path.basename(doc_path)}")
            
            doc_path = os.path.abspath(doc_path)
            output_path = os.path.abspath(output_path)
            
            if self.tools['ms_word'] and HAS_WIN32COM:
                return self._word_to_pdf_msword(doc_path, output_path)
            
            if self.tools['wps'] and HAS_WIN32COM:
                return self._word_to_pdf_wps(doc_path, output_path)
            
            if self.tools['libreoffice']:
                return self._doc_to_pdf_libreoffice(doc_path, output_path)
            
            self.log("❌ 未找到可用的转换工具")
            return False
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
    
    def _word_to_pdf_msword(self, doc_path, output_path):
        try:
            self.log("  使用 Microsoft Word 转换...")
            
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            word.DisplayAlerts = False
            
            try:
                doc = word.Documents.Open(doc_path)
                doc.SaveAs(output_path, FileFormat=17)
                doc.Close()
                self.log(f"✅ PDF保存成功: {output_path}")
                return True
            finally:
                word.Quit()
                
        except Exception as e:
            self.log(f"  Word转换失败: {e}")
            return False
    
    def _word_to_pdf_wps(self, doc_path, output_path):
        try:
            self.log("  使用 WPS 转换...")
            
            try:
                wps = win32com.client.Dispatch("KWPS.Application")
            except:
                wps = win32com.client.Dispatch("KET.Application")
            
            wps.Visible = False
            
            try:
                doc = wps.Documents.Open(doc_path)
                doc.ExportAsFixedFormat(output_path, 17)
                doc.Close()
                self.log(f"✅ PDF保存成功: {output_path}")
                return True
            finally:
                wps.Quit()
                
        except Exception as e:
            self.log(f"  WPS转换失败: {e}")
            return False
    
    def _doc_to_pdf_libreoffice(self, doc_path, output_path):
        try:
            self.log("  使用 LibreOffice 转换...")
            
            soffice = self.tools['libreoffice']
            output_dir = os.path.dirname(output_path)
            
            cmd = [soffice, '--headless', '--convert-to', 'pdf', '--outdir', output_dir, doc_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            expected_output = os.path.join(output_dir, Path(doc_path).stem + '.pdf')
            
            if os.path.exists(expected_output):
                if expected_output != output_path:
                    shutil.move(expected_output, output_path)
                self.log(f"✅ PDF保存成功: {output_path}")
                return True
            else:
                self.log(f"  转换失败: {result.stderr}")
                return False
                
        except Exception as e:
            self.log(f"  LibreOffice转换失败: {e}")
            return False
    
    # ==================== 批量表格转PDF ====================
    def spreadsheets_to_pdf(self, file_paths, output_folder):
        """批量表格转PDF"""
        try:
            self.controller.is_running = True
            total = len(file_paths)
            success_count = 0
            
            self.log(f"🔄 批量转换 {total} 个表格...")
            os.makedirs(output_folder, exist_ok=True)
            
            for i, file_path in enumerate(file_paths):
                if not self.controller.check_pause():
                    break
                
                self.progress(i + 1, total)
                base_name = Path(file_path).stem
                output_path = os.path.join(output_folder, f"{base_name}.pdf")
                
                if self.spreadsheet_to_pdf(file_path, output_path):
                    success_count += 1
            
            self.log(f"✅ 完成！成功 {success_count}/{total}")
            return success_count > 0
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
        finally:
            self.controller.is_running = False
    
    def spreadsheet_to_pdf(self, file_path, output_path):
        """单个表格转PDF"""
        try:
            self.log(f"🔄 转换表格: {os.path.basename(file_path)}")
            
            file_path = os.path.abspath(file_path)
            output_path = os.path.abspath(output_path)
            
            if self.tools['ms_excel'] and HAS_WIN32COM:
                return self._excel_to_pdf_msexcel(file_path, output_path)
            
            if self.tools['wps'] and HAS_WIN32COM:
                return self._excel_to_pdf_wps(file_path, output_path)
            
            if self.tools['libreoffice']:
                return self._doc_to_pdf_libreoffice(file_path, output_path)
            
            self.log("❌ 未找到可用的转换工具")
            return False
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
    
    def _excel_to_pdf_msexcel(self, file_path, output_path):
        try:
            self.log("  使用 Microsoft Excel 转换...")
            
            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = False
            excel.DisplayAlerts = False
            
            try:
                wb = excel.Workbooks.Open(file_path)
                wb.ExportAsFixedFormat(0, output_path)
                wb.Close(False)
                self.log(f"✅ PDF保存成功: {output_path}")
                return True
            finally:
                excel.Quit()
                
        except Exception as e:
            self.log(f"  Excel转换失败: {e}")
            return False
    
    def _excel_to_pdf_wps(self, file_path, output_path):
        try:
            self.log("  使用 WPS 表格转换...")
            
            try:
                et = win32com.client.Dispatch("KET.Application")
            except:
                et = win32com.client.Dispatch("ET.Application")
            
            et.Visible = False
            
            try:
                wb = et.Workbooks.Open(file_path)
                wb.ExportAsFixedFormat(0, output_path)
                wb.Close(False)
                self.log(f"✅ PDF保存成功: {output_path}")
                return True
            finally:
                et.Quit()
                
        except Exception as e:
            self.log(f"  WPS转换失败: {e}")
            return False
    
    # ==================== 批量网页转PDF ====================
    def urls_to_pdf(self, urls, output_folder):
        """批量网页转PDF"""
        try:
            self.controller.is_running = True
            total = len(urls)
            success_count = 0
            
            self.log(f"🔄 批量转换 {total} 个网页...")
            os.makedirs(output_folder, exist_ok=True)
            
            for i, url in enumerate(urls):
                if not self.controller.check_pause():
                    break
                
                self.progress(i + 1, total)
                # 从URL生成文件名
                from urllib.parse import urlparse
                parsed = urlparse(url)
                base_name = parsed.netloc.replace('.', '_') + parsed.path.replace('/', '_')
                base_name = base_name[:50]  # 限制长度
                if not base_name:
                    base_name = f"webpage_{i+1}"
                output_path = os.path.join(output_folder, f"{base_name}.pdf")
                
                if self.url_to_pdf(url, output_path):
                    success_count += 1
            
            self.log(f"✅ 完成！成功 {success_count}/{total}")
            return success_count > 0
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
        finally:
            self.controller.is_running = False
    
    def url_to_pdf(self, url, output_path):
        """单个网页转PDF"""
        try:
            self.log(f"🔄 转换网页: {url}")
            
            if not HAS_SELENIUM:
                self.log("❌ 未安装selenium")
                return False
            
            edge_driver = self.driver_manager.get_edge_driver()
            if edge_driver and self.tools['edge']:
                result = self._url_to_pdf_with_driver(url, output_path, 'edge', edge_driver)
                if result:
                    return True
            
            chrome_driver = self.driver_manager.get_chrome_driver()
            if chrome_driver and self.tools['chrome']:
                result = self._url_to_pdf_with_driver(url, output_path, 'chrome', chrome_driver)
                if result:
                    return True
            
            self.log("  尝试自动下载驱动...")
            result = self._url_to_pdf_auto_driver(url, output_path)
            if result:
                return True
            
            self.log("❌ 未找到浏览器驱动")
            return False
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
    
    def _url_to_pdf_with_driver(self, url, output_path, browser, driver_path):
        driver = None
        try:
            self.log(f"  使用 {browser.upper()} 浏览器渲染...")
            
            if browser == 'edge':
                options = EdgeOptions()
                options.add_argument('--headless=new')
                options.add_argument('--disable-gpu')
                options.add_argument('--no-sandbox')
                options.add_argument('--disable-dev-shm-usage')
                options.add_argument('--window-size=1920,1080')
                options.add_argument('--hide-scrollbars')
                options.add_argument('--disable-extensions')
                options.add_argument('--disable-infobars')
                
                service = EdgeService(executable_path=driver_path)
                driver = webdriver.Edge(service=service, options=options)
            else:
                options = ChromeOptions()
                options.add_argument('--headless=new')
                options.add_argument('--disable-gpu')
                options.add_argument('--no-sandbox')
                options.add_argument('--disable-dev-shm-usage')
                options.add_argument('--window-size=1920,1080')
                options.add_argument('--hide-scrollbars')
                
                service = ChromeService(executable_path=driver_path)
                driver = webdriver.Chrome(service=service, options=options)
            
            return self._capture_webpage_to_pdf(driver, url, output_path)
            
        except Exception as e:
            self.log(f"  {browser}转换失败: {e}")
            return False
        finally:
            if driver:
                try:
                    driver.quit()
                except:
                    pass
    
    def _url_to_pdf_auto_driver(self, url, output_path):
        driver = None
        try:
            try:
                from webdriver_manager.microsoft import EdgeChromiumDriverManager
                from webdriver_manager.chrome import ChromeDriverManager
                HAS_MANAGER = True
            except ImportError:
                HAS_MANAGER = False
                self.log("  未安装 webdriver-manager")
                return False
            
            if self.tools['edge']:
                try:
                    self.log("  下载Edge驱动...")
                    options = EdgeOptions()
                    options.add_argument('--headless=new')
                    options.add_argument('--disable-gpu')
                    options.add_argument('--no-sandbox')
                    options.add_argument('--window-size=1920,1080')
                    
                    service = EdgeService(EdgeChromiumDriverManager().install())
                    driver = webdriver.Edge(service=service, options=options)
                    return self._capture_webpage_to_pdf(driver, url, output_path)
                except Exception as e:
                    self.log(f"  Edge自动驱动失败: {e}")
            
            if self.tools['chrome']:
                try:
                    self.log("  下载Chrome驱动...")
                    options = ChromeOptions()
                    options.add_argument('--headless=new')
                    options.add_argument('--disable-gpu')
                    options.add_argument('--no-sandbox')
                    options.add_argument('--window-size=1920,1080')
                    
                    service = ChromeService(ChromeDriverManager().install())
                    driver = webdriver.Chrome(service=service, options=options)
                    return self._capture_webpage_to_pdf(driver, url, output_path)
                except Exception as e:
                    self.log(f"  Chrome自动驱动失败: {e}")
            
            return False
            
        except Exception as e:
            self.log(f"  自动驱动失败: {e}")
            return False
        finally:
            if driver:
                try:
                    driver.quit()
                except:
                    pass
    
    def _capture_webpage_to_pdf(self, driver, url, output_path):
        try:
            self.log("  加载网页...")
            driver.get(url)
            
            time.sleep(3)
            
            try:
                WebDriverWait(driver, 15).until(
                    EC.presence_of_element_located((By.TAG_NAME, "body"))
                )
            except:
                pass
            
            self.log("  加载完整内容...")
            last_height = driver.execute_script("return document.body.scrollHeight")
            
            scroll_count = 0
            while scroll_count < 30:
                if not self.controller.check_pause():
                    return False
                
                driver.execute_script("window.scrollTo(0, document.body.scrollHeight);")
                time.sleep(0.5)
                
                new_height = driver.execute_script("return document.body.scrollHeight")
                if new_height == last_height:
                    break
                last_height = new_height
                scroll_count += 1
            
            driver.execute_script("window.scrollTo(0, 0);")
            time.sleep(1)
            
            total_height = driver.execute_script("return document.body.scrollHeight")
            self.log(f"  页面高度: {total_height}px")
            
            try:
                self.log("  生成PDF...")
                
                pdf_data = driver.execute_cdp_cmd('Page.printToPDF', {
                    'landscape': False,
                    'displayHeaderFooter': False,
                    'printBackground': True,
                    'preferCSSPageSize': True,
                    'scale': 1,
                    'paperWidth': 8.27,
                    'paperHeight': 11.69,
                    'marginTop': 0.4,
                    'marginBottom': 0.4,
                    'marginLeft': 0.4,
                    'marginRight': 0.4,
                })
                
                import base64
                pdf_bytes = base64.b64decode(pdf_data['data'])
                
                with open(output_path, 'wb') as f:
                    f.write(pdf_bytes)
                
                self.log(f"✅ PDF保存成功: {output_path}")
                return True
                
            except Exception as e:
                self.log(f"  CDP打印失败: {e}")
                return self._capture_screenshot_to_pdf(driver, output_path, total_height)
            
        except Exception as e:
            self.log(f"  网页捕获失败: {e}")
            return False
    
    def _capture_screenshot_to_pdf(self, driver, output_path, total_height):
        try:
            driver.set_window_size(1920, min(total_height + 200, 16000))
            time.sleep(1)
            
            screenshot = driver.get_screenshot_as_png()
            img = Image.open(BytesIO(screenshot))
            
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            img.save(output_path, 'PDF', resolution=150)
            img.close()
            
            self.log(f"✅ PDF保存成功: {output_path}")
            return True
            
        except Exception as e:
            self.log(f"  截图失败: {e}")
            return False
    
    # ==================== 图片转PDF ====================
    def images_to_pdf(self, image_paths, output_path, quality='high'):
        try:
            self.controller.is_running = True
            total = len(image_paths)
            self.log(f"🔄 转换 {total} 张图片为PDF...")
            
            quality_settings = {
                'high': {'dpi': 300, 'quality': 100},
                'medium': {'dpi': 150, 'quality': 85},
                'low': {'dpi': 72, 'quality': 70}
            }
            settings = quality_settings.get(quality, quality_settings['high'])
            
            images = []
            first_image = None
            
            for i, img_path in enumerate(image_paths):
                if not self.controller.check_pause():
                    return False
                
                self.log(f"  处理 {i+1}/{total}: {os.path.basename(img_path)}")
                self.progress(i + 1, total)
                
                img = Image.open(img_path)
                
                if img.mode == 'RGBA':
                    bg = Image.new('RGB', img.size, (255, 255, 255))
                    bg.paste(img, mask=img.split()[3])
                    img = bg
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                
                if first_image is None:
                    first_image = img
                else:
                    images.append(img)
            
            if first_image and not self.controller.should_stop():
                first_image.save(
                    output_path, 'PDF',
                    resolution=settings['dpi'],
                    save_all=True,
                    append_images=images,
                    quality=settings['quality']
                )
                self.log(f"✅ PDF保存成功: {output_path}")
                return True
            
            return False
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
        finally:
            self.controller.is_running = False
    
    # ==================== 图片转PPT ====================
    def images_to_ppt(self, image_paths, output_path, quality='high'):
        try:
            self.controller.is_running = True
            total = len(image_paths)
            self.log(f"🔄 转换 {total} 张图片为PPT...")
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]
            
            for i, img_path in enumerate(image_paths):
                if not self.controller.check_pause():
                    return False
                
                self.log(f"  处理 {i+1}/{total}: {os.path.basename(img_path)}")
                self.progress(i + 1, total)
                
                slide = prs.slides.add_slide(blank_layout)
                img = Image.open(img_path)
                
                temp_path = img_path
                if quality != 'high' or img.mode == 'RGBA':
                    if quality != 'high':
                        max_size = (1920, 1080) if quality == 'medium' else (1280, 720)
                        img.thumbnail(max_size, Image.Resampling.LANCZOS)
                    
                    if img.mode == 'RGBA':
                        bg = Image.new('RGB', img.size, (255, 255, 255))
                        bg.paste(img, mask=img.split()[3])
                        img = bg
                    
                    temp_path = tempfile.mktemp(suffix='.jpg')
                    img.save(temp_path, 'JPEG', quality=95 if quality == 'high' else 85)
                
                img_w, img_h = img.size
                slide_w = prs.slide_width
                slide_h = prs.slide_height
                
                img_w_emu = Emu(img_w * 914400 / 96)
                img_h_emu = Emu(img_h * 914400 / 96)
                
                ratio = min(slide_w / img_w_emu, slide_h / img_h_emu) * 0.95
                new_w = int(img_w_emu * ratio)
                new_h = int(img_h_emu * ratio)
                
                left = (slide_w - new_w) // 2
                top = (slide_h - new_h) // 2
                
                slide.shapes.add_picture(temp_path, left, top, new_w, new_h)
                
                if temp_path != img_path and os.path.exists(temp_path):
                    try:
                        os.remove(temp_path)
                    except:
                        pass
            
            if not self.controller.should_stop():
                prs.save(output_path)
                self.log(f"✅ PPT保存成功: {output_path}")
                return True
            
            return False
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
        finally:
            self.controller.is_running = False
    
    # ==================== 批量PDF转PPT ====================
    def pdfs_to_ppt(self, pdf_paths, output_folder, dpi=150):
        """批量PDF转PPT"""
        try:
            self.controller.is_running = True
            total = len(pdf_paths)
            success_count = 0
            
            self.log(f"🔄 批量转换 {total} 个PDF...")
            os.makedirs(output_folder, exist_ok=True)
            
            for i, pdf_path in enumerate(pdf_paths):
                if not self.controller.check_pause():
                    break
                
                base_name = Path(pdf_path).stem
                output_path = os.path.join(output_folder, f"{base_name}.pptx")
                
                # 重置控制器状态用于子任务
                self.controller.is_running = True
                
                if self._pdf_to_ppt_single(pdf_path, output_path, dpi, i + 1, total):
                    success_count += 1
            
            self.log(f"✅ 完成！成功 {success_count}/{total}")
            return success_count > 0
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
        finally:
            self.controller.is_running = False
    
    def _pdf_to_ppt_single(self, pdf_path, output_path, dpi, current_file, total_files):
        """单个PDF转PPT"""
        try:
            self.log(f"🔄 [{current_file}/{total_files}] 转换: {os.path.basename(pdf_path)}")
            
            pdf_doc = fitz.open(pdf_path)
            total = len(pdf_doc)
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]
            
            for page_num in range(total):
                if not self.controller.check_pause():
                    pdf_doc.close()
                    return False
                
                self.log(f"  处理页面 {page_num + 1}/{total}")
                self.progress(page_num + 1, total)
                
                page = pdf_doc[page_num]
                mat = fitz.Matrix(dpi/72, dpi/72)
                pix = page.get_pixmap(matrix=mat)
                
                temp_img = tempfile.mktemp(suffix='.png')
                pix.save(temp_img)
                
                slide = prs.slides.add_slide(blank_layout)
                
                img = Image.open(temp_img)
                img_w, img_h = img.size
                img.close()
                
                img_w_emu = Emu(img_w * 914400 / dpi)
                img_h_emu = Emu(img_h * 914400 / dpi)
                
                ratio = min(prs.slide_width / img_w_emu, prs.slide_height / img_h_emu) * 0.95
                new_w = int(img_w_emu * ratio)
                new_h = int(img_h_emu * ratio)
                
                left = (prs.slide_width - new_w) // 2
                top = (prs.slide_height - new_h) // 2
                
                slide.shapes.add_picture(temp_img, left, top, new_w, new_h)
                
                try:
                    os.remove(temp_img)
                except:
                    pass
                
                pix = None
            
            pdf_doc.close()
            
            if not self.controller.should_stop():
                prs.save(output_path)
                self.log(f"✅ PPT保存成功: {output_path}")
                return True
            
            return False
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
    
    def pdf_to_ppt(self, pdf_path, output_path, dpi=150):
        """单个PDF转PPT（保持兼容性）"""
        try:
            self.controller.is_running = True
            return self._pdf_to_ppt_single(pdf_path, output_path, dpi, 1, 1)
        finally:
            self.controller.is_running = False
    
    # ==================== 批量PDF转图片 ====================
    def pdfs_to_images(self, pdf_paths, output_folder, dpi=200, img_format='png'):
        """批量PDF转图片"""
        try:
            self.controller.is_running = True
            total = len(pdf_paths)
            success_count = 0
            
            self.log(f"🔄 批量转换 {total} 个PDF为图片...")
            os.makedirs(output_folder, exist_ok=True)
            
            for i, pdf_path in enumerate(pdf_paths):
                if not self.controller.check_pause():
                    break
                
                base_name = Path(pdf_path).stem
                pdf_output_folder = os.path.join(output_folder, base_name)
                
                self.controller.is_running = True
                
                if self._pdf_to_images_single(pdf_path, pdf_output_folder, dpi, img_format, i + 1, total):
                    success_count += 1
            
            self.log(f"✅ 完成！成功 {success_count}/{total}")
            return success_count > 0
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
        finally:
            self.controller.is_running = False
    
    def _pdf_to_images_single(self, pdf_path, output_folder, dpi, img_format, current_file, total_files):
        """单个PDF转图片"""
        try:
            self.log(f"🔄 [{current_file}/{total_files}] 转换: {os.path.basename(pdf_path)}")
            
            os.makedirs(output_folder, exist_ok=True)
            
            pdf_doc = fitz.open(pdf_path)
            total = len(pdf_doc)
            base_name = Path(pdf_path).stem
            
            self.log(f"  共 {total} 页，DPI: {dpi}")
            
            for page_num in range(total):
                if not self.controller.check_pause():
                    pdf_doc.close()
                    return False
                
                self.log(f"  处理页面 {page_num + 1}/{total}")
                self.progress(page_num + 1, total)
                
                page = pdf_doc[page_num]
                mat = fitz.Matrix(dpi/72, dpi/72)
                pix = page.get_pixmap(matrix=mat)
                
                out_path = os.path.join(output_folder, f"{base_name}_page_{page_num + 1:03d}.{img_format}")
                
                if img_format.lower() in ['jpg', 'jpeg']:
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img.save(out_path, "JPEG", quality=95)
                    img.close()
                else:
                    pix.save(out_path)
                
                pix = None
            
            pdf_doc.close()
            self.log(f"✅ 共 {total} 张图片保存到: {output_folder}")
            return True
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
    
    def pdf_to_images(self, pdf_path, output_folder, dpi=200, img_format='png'):
        """单个PDF转图片（保持兼容性）"""
        try:
            self.controller.is_running = True
            return self._pdf_to_images_single(pdf_path, output_folder, dpi, img_format, 1, 1)
        finally:
            self.controller.is_running = False
    
    # ==================== 批量提取PDF图片 ====================
    def extract_images_from_pdfs(self, pdf_paths, output_folder):
        """批量提取PDF中的图片"""
        try:
            self.controller.is_running = True
            total = len(pdf_paths)
            success_count = 0
            
            self.log(f"🔄 批量提取 {total} 个PDF中的图片...")
            os.makedirs(output_folder, exist_ok=True)
            
            for i, pdf_path in enumerate(pdf_paths):
                if not self.controller.check_pause():
                    break
                
                base_name = Path(pdf_path).stem
                pdf_output_folder = os.path.join(output_folder, base_name)
                
                self.controller.is_running = True
                
                if self._extract_images_single(pdf_path, pdf_output_folder, i + 1, total):
                    success_count += 1
            
            self.log(f"✅ 完成！成功 {success_count}/{total}")
            return success_count > 0
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
        finally:
            self.controller.is_running = False
    
    def _extract_images_single(self, pdf_path, output_folder, current_file, total_files):
        """单个PDF提取图片"""
        try:
            self.log(f"🔄 [{current_file}/{total_files}] 提取: {os.path.basename(pdf_path)}")
            
            os.makedirs(output_folder, exist_ok=True)
            
            pdf_doc = fitz.open(pdf_path)
            total = len(pdf_doc)
            image_count = 0
            
            for page_num in range(total):
                if not self.controller.check_pause():
                    pdf_doc.close()
                    return False
                
                page = pdf_doc[page_num]
                image_list = page.get_images(full=True)
                
                self.log(f"  页面 {page_num + 1}: {len(image_list)} 张图片")
                self.progress(page_num + 1, total)
                
                for img_idx, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = pdf_doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        image_count += 1
                        out_path = os.path.join(output_folder, f"image_page{page_num + 1}_{img_idx + 1}.{image_ext}")
                        
                        with open(out_path, "wb") as f:
                            f.write(image_bytes)
                    except:
                        continue
            
            pdf_doc.close()
            
            if image_count == 0:
                self.log("⚠️ PDF中没有找到图片")
            else:
                self.log(f"✅ 共提取 {image_count} 张图片")
            
            return True
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
    
    def extract_images_from_pdf(self, pdf_path, output_folder):
        """单个PDF提取图片（保持兼容性）"""
        try:
            self.controller.is_running = True
            return self._extract_images_single(pdf_path, output_folder, 1, 1)
        finally:
            self.controller.is_running = False
    
    # ==================== 图片转WebP ====================
    def images_to_webp(self, input_paths, output_folder, quality=85, resize_percent=100):
        try:
            self.controller.is_running = True
            total = len(input_paths)
            self.log(f"🔄 转换 {total} 张图片为WebP...")
            
            os.makedirs(output_folder, exist_ok=True)
            success_count = 0
            
            for i, img_path in enumerate(input_paths):
                if not self.controller.check_pause():
                    return False
                
                try:
                    self.log(f"  处理 {i+1}/{total}: {os.path.basename(img_path)}")
                    self.progress(i + 1, total)
                    
                    img = Image.open(img_path)
                    
                    if resize_percent != 100:
                        new_w = int(img.width * resize_percent / 100)
                        new_h = int(img.height * resize_percent / 100)
                        img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
                    
                    base_name = Path(img_path).stem
                    out_path = os.path.join(output_folder, f"{base_name}.webp")
                    
                    img.save(out_path, 'WEBP', quality=quality, method=6)
                    img.close()
                    
                    success_count += 1
                    
                except Exception as e:
                    self.log(f"    ⚠️ 失败: {e}")
            
            self.log(f"✅ 成功转换 {success_count}/{total} 张图片")
            return True
            
        except Exception as e:
            self.log(f"❌ 错误: {str(e)}")
            return False
        finally:
            self.controller.is_running = False
    
    def folder_to_webp(self, input_folder, output_folder, quality=85, resize_percent=100):
        extensions = {'.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff'}
        files = []
        
        for root, dirs, filenames in os.walk(input_folder):
            for f in filenames:
                if Path(f).suffix.lower() in extensions:
                    files.append(os.path.join(root, f))
        
        if not files:
            self.log("⚠️ 文件夹中没有图片")
            return False
        
        self.log(f"📁 找到 {len(files)} 张图片")
        return self.images_to_webp(files, output_folder, quality, resize_percent)


# ==================== GUI界面 ====================
class ConverterGUI:
    """GUI界面"""
    
    def __init__(self, root):
        self.root = root
        self.root.title("📄 多功能文档转换工具 v4.2 - 支持批量多文件")
        self.root.geometry("1000x850")
        self.root.minsize(900, 700)
        
        self.style = ttk.Style()
        try:
            self.style.theme_use('vista')
        except:
            pass
        
        self.converter = DocumentConverter(self.log_message, self.update_progress)
        
        # 文件列表存储
        self.selected_files = []      # 图片文件
        self.doc_files = []           # 文档文件
        self.sheet_files = []         # 表格文件
        self.pdf_files = []           # PDF文件
        self.url_list = []            # URL列表
        
        self.create_widgets()
    
    def create_widgets(self):
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 标题
        title_frame = ttk.Frame(main_frame)
        title_frame.pack(fill=tk.X, pady=(0, 10))
        
        ttk.Label(title_frame, text="🔄 多功能文档转换工具 v4.2", font=('微软雅黑', 16, 'bold')).pack()
        ttk.Label(title_frame, text="✨ 支持批量多文件选择 | 完整保留原始样式", font=('微软雅黑', 9), foreground='green').pack()
        
        # Notebook
        self.notebook = ttk.Notebook(main_frame)
        self.notebook.pack(fill=tk.BOTH, expand=True, pady=(5, 10))
        
        self.create_image_tab()
        self.create_document_tab()
        self.create_spreadsheet_tab()
        self.create_webpage_tab()
        self.create_pdf_tab()
        self.create_webp_tab()
        self.create_help_tab()
        
        # 进度条
        progress_frame = ttk.Frame(main_frame)
        progress_frame.pack(fill=tk.X, pady=(0, 5))
        
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(progress_frame, variable=self.progress_var, maximum=100)
        self.progress_bar.pack(fill=tk.X, side=tk.LEFT, expand=True, padx=(0, 10))
        
        self.progress_label = ttk.Label(progress_frame, text="就绪", width=15)
        self.progress_label.pack(side=tk.RIGHT)
        
        # 控制按钮
        control_frame = ttk.Frame(main_frame)
        control_frame.pack(fill=tk.X, pady=(0, 5))
        
        ttk.Button(control_frame, text="⏸️ 暂停", command=self.pause_task, width=12).pack(side=tk.LEFT, padx=2)
        ttk.Button(control_frame, text="▶️ 继续", command=self.resume_task, width=12).pack(side=tk.LEFT, padx=2)
        ttk.Button(control_frame, text="⏹️ 终止", command=self.stop_task, width=12).pack(side=tk.LEFT, padx=2)
        ttk.Button(control_frame, text="🗑️ 清除日志", command=self.clear_log, width=12).pack(side=tk.RIGHT, padx=2)
        
        # 日志
        log_frame = ttk.LabelFrame(main_frame, text="📋 操作日志", padding="5")
        log_frame.pack(fill=tk.BOTH, expand=True)
        
        self.log_text = ScrolledText(log_frame, height=8, font=('Consolas', 9), bg='#1e1e1e', fg='#d4d4d4')
        self.log_text.pack(fill=tk.BOTH, expand=True)
        
        self.log_message("✅ 程序启动成功！支持批量多文件选择")
        self.log_message("=" * 60)
    
    def create_file_list_widget(self, parent, list_var_name, file_types, title="文件列表"):
        """创建通用的文件列表组件"""
        frame = ttk.LabelFrame(parent, text=f"📁 {title}", padding="10")
        frame.pack(fill=tk.BOTH, expand=True, pady=(0, 10))
        
        btn_frame = ttk.Frame(frame)
        btn_frame.pack(fill=tk.X, pady=(0, 5))
        
        add_btn = ttk.Button(btn_frame, text="➕ 添加文件", width=12)
        add_btn.pack(side=tk.LEFT, padx=2)
        
        clear_btn = ttk.Button(btn_frame, text="🗑️ 清除全部", width=12)
        clear_btn.pack(side=tk.LEFT, padx=2)
        
        del_btn = ttk.Button(btn_frame, text="❌ 删除选中", width=12)
        del_btn.pack(side=tk.LEFT, padx=2)
        
        count_label = ttk.Label(btn_frame, text="已选择: 0 个文件")
        count_label.pack(side=tk.RIGHT, padx=5)
        
        list_frame = ttk.Frame(frame)
        list_frame.pack(fill=tk.BOTH, expand=True)
        
        scrollbar = ttk.Scrollbar(list_frame)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        listbox = tk.Listbox(list_frame, height=5, selectmode=tk.EXTENDED, 
                            yscrollcommand=scrollbar.set, font=('Consolas', 9))
        listbox.pack(fill=tk.BOTH, expand=True)
        scrollbar.config(command=listbox.yview)
        
        return frame, listbox, add_btn, clear_btn, del_btn, count_label
    
    def create_image_tab(self):
        tab = ttk.Frame(self.notebook, padding="15")
        self.notebook.add(tab, text="  🖼️ 图片转换  ")
        
        # 文件列表
        _, self.image_listbox, add_btn, clear_btn, del_btn, self.image_count_label = \
            self.create_file_list_widget(tab, 'selected_files', 
                [("图片", "*.png *.jpg *.jpeg *.bmp *.gif *.tiff *.webp")], "图片文件")
        
        add_btn.config(command=self.add_images)
        clear_btn.config(command=self.clear_images)
        del_btn.config(command=self.delete_selected_images)
        
        # 排序按钮
        sort_frame = ttk.Frame(tab)
        sort_frame.pack(fill=tk.X, pady=(0, 10))
        ttk.Button(sort_frame, text="⬆️ 上移", command=self.move_up, width=8).pack(side=tk.LEFT, padx=2)
        ttk.Button(sort_frame, text="⬇️ 下移", command=self.move_down, width=8).pack(side=tk.LEFT, padx=2)
        
        # 质量设置
        quality_frame = ttk.LabelFrame(tab, text="⚙️ 质量设置", padding="10")
        quality_frame.pack(fill=tk.X, pady=(0, 10))
        
        self.image_quality = tk.StringVar(value='high')
        for val, text in [('high', '高 (300 DPI)'), ('medium', '中 (150 DPI)'), ('low', '低 (72 DPI)')]:
            ttk.Radiobutton(quality_frame, text=text, variable=self.image_quality, value=val).pack(side=tk.LEFT, padx=10)
        
        # 转换按钮
        convert_frame = ttk.Frame(tab)
        convert_frame.pack(fill=tk.X)
        
        ttk.Button(convert_frame, text="📄 转为 PDF", command=lambda: self.run_task(self.convert_images_to_pdf), width=15).pack(side=tk.LEFT, padx=5)
        ttk.Button(convert_frame, text="📊 转为 PPT", command=lambda: self.run_task(self.convert_images_to_ppt), width=15).pack(side=tk.LEFT, padx=5)
    
    def create_document_tab(self):
        tab = ttk.Frame(self.notebook, padding="15")
        self.notebook.add(tab, text="  📝 文档转换  ")
        
        # 提示
        info_frame = ttk.LabelFrame(tab, text="💡 提示", padding="10")
        info_frame.pack(fill=tk.X, pady=(0, 10))
        ttk.Label(info_frame, text="支持批量选择多个Word文档，使用 Office / WPS / LibreOffice 转换", 
                 font=('微软雅黑', 9)).pack(anchor=tk.W)
        
        # 文件列表
        _, self.doc_listbox, add_btn, clear_btn, del_btn, self.doc_count_label = \
            self.create_file_list_widget(tab, 'doc_files', 
                [("Word文档", "*.docx *.doc *.wps *.rtf")], "Word文档")
        
        add_btn.config(command=self.add_documents)
        clear_btn.config(command=self.clear_documents)
        del_btn.config(command=self.delete_selected_documents)
        
        # 转换按钮
        convert_frame = ttk.Frame(tab)
        convert_frame.pack(fill=tk.X, pady=10)
        ttk.Button(convert_frame, text="📄 批量转为 PDF", 
                  command=lambda: self.run_task(self.convert_docs_to_pdf), width=20).pack(side=tk.LEFT, padx=5)
    
    def create_spreadsheet_tab(self):
        tab = ttk.Frame(self.notebook, padding="15")
        self.notebook.add(tab, text="  📊 表格转换  ")
        
        # 提示
        info_frame = ttk.LabelFrame(tab, text="💡 提示", padding="10")
        info_frame.pack(fill=tk.X, pady=(0, 10))
        ttk.Label(info_frame, text="支持批量选择多个Excel表格，使用 Office / WPS / LibreOffice 转换", 
                 font=('微软雅黑', 9)).pack(anchor=tk.W)
        
        # 文件列表
        _, self.sheet_listbox, add_btn, clear_btn, del_btn, self.sheet_count_label = \
            self.create_file_list_widget(tab, 'sheet_files', 
                [("Excel表格", "*.xlsx *.xls *.csv")], "Excel表格")
        
        add_btn.config(command=self.add_spreadsheets)
        clear_btn.config(command=self.clear_spreadsheets)
        del_btn.config(command=self.delete_selected_spreadsheets)
        
        # 转换按钮
        convert_frame = ttk.Frame(tab)
        convert_frame.pack(fill=tk.X, pady=10)
        ttk.Button(convert_frame, text="📄 批量转为 PDF", 
                  command=lambda: self.run_task(self.convert_sheets_to_pdf), width=20).pack(side=tk.LEFT, padx=5)
    
    def create_webpage_tab(self):
        tab = ttk.Frame(self.notebook, padding="15")
        self.notebook.add(tab, text="  🌐 网页转换  ")
        
        # 提示
        info_frame = ttk.LabelFrame(tab, text="💡 提示", padding="10")
        info_frame.pack(fill=tk.X, pady=(0, 10))
        ttk.Label(info_frame, text="支持批量转换多个网页URL，每行一个URL", 
                 font=('微软雅黑', 9)).pack(anchor=tk.W)
        
        # URL输入
        url_frame = ttk.LabelFrame(tab, text="🔗 网页地址（每行一个）", padding="10")
        url_frame.pack(fill=tk.BOTH, expand=True, pady=(0, 10))
        
        self.url_text = ScrolledText(url_frame, height=8, font=('Consolas', 10))
        self.url_text.pack(fill=tk.BOTH, expand=True)
        self.url_text.insert(tk.END, "https://example.com\n")
        
        # 按钮
        btn_frame = ttk.Frame(tab)
        btn_frame.pack(fill=tk.X)
        ttk.Button(btn_frame, text="🌐 批量网页 → PDF", 
                  command=lambda: self.run_task(self.convert_urls_to_pdf), width=20).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="🗑️ 清除", 
                  command=lambda: self.url_text.delete(1.0, tk.END), width=10).pack(side=tk.LEFT, padx=5)
    
    def create_pdf_tab(self):
        tab = ttk.Frame(self.notebook, padding="15")
        self.notebook.add(tab, text="  📄 PDF转换  ")
        
        # 文件列表
        _, self.pdf_listbox, add_btn, clear_btn, del_btn, self.pdf_count_label = \
            self.create_file_list_widget(tab, 'pdf_files', 
                [("PDF文件", "*.pdf")], "PDF文件")
        
        add_btn.config(command=self.add_pdfs)
        clear_btn.config(command=self.clear_pdfs)
        del_btn.config(command=self.delete_selected_pdfs)
        
        # 设置
        settings_frame = ttk.LabelFrame(tab, text="⚙️ 设置", padding="10")
        settings_frame.pack(fill=tk.X, pady=(0, 10))
        
        settings_row = ttk.Frame(settings_frame)
        settings_row.pack(fill=tk.X)
        
        ttk.Label(settings_row, text="DPI:").pack(side=tk.LEFT, padx=(0, 5))
        self.pdf_dpi = tk.IntVar(value=150)
        ttk.Spinbox(settings_row, from_=72, to=600, textvariable=self.pdf_dpi, width=8).pack(side=tk.LEFT, padx=(0, 20))
        
        ttk.Label(settings_row, text="格式:").pack(side=tk.LEFT, padx=(0, 5))
        self.image_format = tk.StringVar(value='png')
        ttk.Radiobutton(settings_row, text="PNG", variable=self.image_format, value='png').pack(side=tk.LEFT, padx=5)
        ttk.Radiobutton(settings_row, text="JPG", variable=self.image_format, value='jpg').pack(side=tk.LEFT, padx=5)
        
        # 转换按钮
        convert_frame = ttk.LabelFrame(tab, text="🔄 转换操作", padding="10")
        convert_frame.pack(fill=tk.X)
        
        btn_row = ttk.Frame(convert_frame)
        btn_row.pack(fill=tk.X)
        
        ttk.Button(btn_row, text="📊 批量PDF→PPT", 
                  command=lambda: self.run_task(self.convert_pdfs_to_ppt), width=15).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_row, text="🖼️ 批量PDF→图片", 
                  command=lambda: self.run_task(self.convert_pdfs_to_images), width=15).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_row, text="📤 批量提取图片", 
                  command=lambda: self.run_task(self.extract_pdfs_images), width=15).pack(side=tk.LEFT, padx=5)
    
    def create_webp_tab(self):
        tab = ttk.Frame(self.notebook, padding="15")
        self.notebook.add(tab, text="  🖼️ WebP  ")
        
        # 模式选择
        mode_frame = ttk.LabelFrame(tab, text="📁 输入模式", padding="10")
        mode_frame.pack(fill=tk.X, pady=(0, 10))
        
        self.webp_mode = tk.StringVar(value='files')
        ttk.Radiobutton(mode_frame, text="选择图片文件", variable=self.webp_mode, value='files').pack(side=tk.LEFT, padx=10)
        ttk.Radiobutton(mode_frame, text="选择文件夹", variable=self.webp_mode, value='folder').pack(side=tk.LEFT, padx=10)
        
        # 输入
        input_frame = ttk.LabelFrame(tab, text="📂 输入", padding="10")
        input_frame.pack(fill=tk.X, pady=(0, 10))
        
        self.webp_input = tk.StringVar()
        self.webp_files = []
        
        path_frame = ttk.Frame(input_frame)
        path_frame.pack(fill=tk.X)
        ttk.Entry(path_frame, textvariable=self.webp_input, font=('Consolas', 10)).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))
        ttk.Button(path_frame, text="浏览", command=self.browse_webp_input, width=8).pack(side=tk.RIGHT)
        
        self.webp_count_label = ttk.Label(input_frame, text="未选择")
        self.webp_count_label.pack(anchor=tk.W, pady=(5, 0))
        
        # 设置
        settings_frame = ttk.LabelFrame(tab, text="⚙️ 设置", padding="10")
        settings_frame.pack(fill=tk.X, pady=(0, 10))
        
        quality_row = ttk.Frame(settings_frame)
        quality_row.pack(fill=tk.X, pady=5)
        ttk.Label(quality_row, text="质量:").pack(side=tk.LEFT, padx=(0, 10))
        self.webp_quality = tk.IntVar(value=85)
        quality_scale = ttk.Scale(quality_row, from_=1, to=100, variable=self.webp_quality, orient=tk.HORIZONTAL, length=150)
        quality_scale.pack(side=tk.LEFT, padx=(0, 10))
        self.quality_label = ttk.Label(quality_row, text="85")
        self.quality_label.pack(side=tk.LEFT)
        quality_scale.configure(command=lambda v: self.quality_label.configure(text=str(int(float(v)))))
        
        size_row = ttk.Frame(settings_frame)
        size_row.pack(fill=tk.X, pady=5)
        ttk.Label(size_row, text="尺寸:").pack(side=tk.LEFT, padx=(0, 10))
        self.webp_resize = tk.IntVar(value=100)
        for val, text in [(100, '100%'), (75, '75%'), (50, '50%')]:
            ttk.Radiobutton(size_row, text=text, variable=self.webp_resize, value=val).pack(side=tk.LEFT, padx=5)
        
        ttk.Button(tab, text="🔄 转换为 WebP", command=lambda: self.run_task(self.convert_to_webp), width=20).pack(pady=10)
    
    def create_help_tab(self):
        tab = ttk.Frame(self.notebook, padding="15")
        self.notebook.add(tab, text="  ❓ 帮助  ")
        
        driver_frame = ttk.LabelFrame(tab, text="🔧 浏览器驱动下载", padding="10")
        driver_frame.pack(fill=tk.BOTH, expand=True)
        
        help_text = ScrolledText(driver_frame, height=15, font=('Consolas', 10))
        help_text.pack(fill=tk.BOTH, expand=True)
        
        instructions = self.converter.get_driver_instructions()
        help_text.insert(tk.END, instructions)
        help_text.config(state=tk.DISABLED)
        
        btn_frame = ttk.Frame(driver_frame)
        btn_frame.pack(fill=tk.X, pady=(10, 0))
        
        ttk.Button(btn_frame, text="🔄 重新检测驱动", command=self.refresh_drivers).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="📂 打开程序目录", command=lambda: os.startfile(os.getcwd())).pack(side=tk.LEFT, padx=5)
    
    def refresh_drivers(self):
        self.converter.driver_manager._find_drivers()
        self.log_message("🔄 已重新检测驱动")
    
    # ================ 辅助方法 ================
    
    def log_message(self, message):
        def _log():
            self.log_text.insert(tk.END, message + "\n")
            self.log_text.see(tk.END)
        self.root.after(0, _log)
    
    def update_progress(self, current, total):
        def _update():
            percent = (current / total) * 100 if total > 0 else 0
            self.progress_var.set(percent)
            self.progress_label.config(text=f"{current}/{total} ({percent:.0f}%)")
        self.root.after(0, _update)
    
    def clear_log(self):
        self.log_text.delete(1.0, tk.END)
    
    def run_task(self, func):
        if self.converter.controller.is_running:
            messagebox.showwarning("提示", "有任务正在运行")
            return
        self.converter.controller.reset()
        self.progress_var.set(0)
        self.progress_label.config(text="处理中...")
        threading.Thread(target=func, daemon=True).start()
    
    def pause_task(self):
        if self.converter.controller.is_running:
            self.converter.controller.pause()
            self.log_message("⏸️ 已暂停")
    
    def resume_task(self):
        self.converter.controller.resume()
        self.log_message("▶️ 继续")
    
    def stop_task(self):
        if self.converter.controller.is_running:
            self.converter.controller.stop()
            self.log_message("⏹️ 终止中...")
    
    # ================ 图片文件操作 ================
    
    def add_images(self):
        files = filedialog.askopenfilenames(
            title="选择图片（可多选）", 
            filetypes=[("图片", "*.png *.jpg *.jpeg *.bmp *.gif *.tiff *.webp"), ("所有", "*.*")]
        )
        for f in files:
            if f not in self.selected_files:
                self.selected_files.append(f)
                self.image_listbox.insert(tk.END, os.path.basename(f))
        self._update_image_count()
        if files:
            self.log_message(f"📁 添加 {len(files)} 张图片，共 {len(self.selected_files)} 张")
    
    def clear_images(self):
        self.selected_files.clear()
        self.image_listbox.delete(0, tk.END)
        self._update_image_count()
    
    def delete_selected_images(self):
        sel = list(self.image_listbox.curselection())
        for i in reversed(sel):
            del self.selected_files[i]
            self.image_listbox.delete(i)
        self._update_image_count()
    
    def _update_image_count(self):
        self.image_count_label.config(text=f"已选择: {len(self.selected_files)} 个文件")
    
    def move_up(self):
        sel = self.image_listbox.curselection()
        if not sel or sel[0] == 0:
            return
        for i in sel:
            self.selected_files[i], self.selected_files[i-1] = self.selected_files[i-1], self.selected_files[i]
            text = self.image_listbox.get(i)
            self.image_listbox.delete(i)
            self.image_listbox.insert(i-1, text)
            self.image_listbox.selection_set(i-1)
    
    def move_down(self):
        sel = self.image_listbox.curselection()
        if not sel or sel[-1] >= len(self.selected_files) - 1:
            return
        for i in reversed(sel):
            self.selected_files[i], self.selected_files[i+1] = self.selected_files[i+1], self.selected_files[i]
            text = self.image_listbox.get(i)
            self.image_listbox.delete(i)
            self.image_listbox.insert(i+1, text)
            self.image_listbox.selection_set(i+1)
    
    # ================ 文档文件操作 ================
    
    def add_documents(self):
        files = filedialog.askopenfilenames(
            title="选择文档（可多选）", 
            filetypes=[("Word文档", "*.docx *.doc *.wps *.rtf"), ("所有", "*.*")]
        )
        for f in files:
            if f not in self.doc_files:
                self.doc_files.append(f)
                self.doc_listbox.insert(tk.END, os.path.basename(f))
        self._update_doc_count()
        if files:
            self.log_message(f"📄 添加 {len(files)} 个文档，共 {len(self.doc_files)} 个")
    
    def clear_documents(self):
        self.doc_files.clear()
        self.doc_listbox.delete(0, tk.END)
        self._update_doc_count()
    
    def delete_selected_documents(self):
        sel = list(self.doc_listbox.curselection())
        for i in reversed(sel):
            del self.doc_files[i]
            self.doc_listbox.delete(i)
        self._update_doc_count()
    
    def _update_doc_count(self):
        self.doc_count_label.config(text=f"已选择: {len(self.doc_files)} 个文件")
    
    # ================ 表格文件操作 ================
    
    def add_spreadsheets(self):
        files = filedialog.askopenfilenames(
            title="选择表格（可多选）", 
            filetypes=[("Excel表格", "*.xlsx *.xls *.csv"), ("所有", "*.*")]
        )
        for f in files:
            if f not in self.sheet_files:
                self.sheet_files.append(f)
                self.sheet_listbox.insert(tk.END, os.path.basename(f))
        self._update_sheet_count()
        if files:
            self.log_message(f"📊 添加 {len(files)} 个表格，共 {len(self.sheet_files)} 个")
    
    def clear_spreadsheets(self):
        self.sheet_files.clear()
        self.sheet_listbox.delete(0, tk.END)
        self._update_sheet_count()
    
    def delete_selected_spreadsheets(self):
        sel = list(self.sheet_listbox.curselection())
        for i in reversed(sel):
            del self.sheet_files[i]
            self.sheet_listbox.delete(i)
        self._update_sheet_count()
    
    def _update_sheet_count(self):
        self.sheet_count_label.config(text=f"已选择: {len(self.sheet_files)} 个文件")
    
    # ================ PDF文件操作 ================
    
    def add_pdfs(self):
        files = filedialog.askopenfilenames(
            title="选择PDF（可多选）", 
            filetypes=[("PDF文件", "*.pdf"), ("所有", "*.*")]
        )
        for f in files:
            if f not in self.pdf_files:
                self.pdf_files.append(f)
                self.pdf_listbox.insert(tk.END, os.path.basename(f))
        self._update_pdf_count()
        if files:
            self.log_message(f"📄 添加 {len(files)} 个PDF，共 {len(self.pdf_files)} 个")
    
    def clear_pdfs(self):
        self.pdf_files.clear()
        self.pdf_listbox.delete(0, tk.END)
        self._update_pdf_count()
    
    def delete_selected_pdfs(self):
        sel = list(self.pdf_listbox.curselection())
        for i in reversed(sel):
            del self.pdf_files[i]
            self.pdf_listbox.delete(i)
        self._update_pdf_count()
    
    def _update_pdf_count(self):
        self.pdf_count_label.config(text=f"已选择: {len(self.pdf_files)} 个文件")
    
    # ================ WebP操作 ================
    
    def browse_webp_input(self):
        mode = self.webp_mode.get()
        if mode == 'files':
            files = filedialog.askopenfilenames(
                title="选择图片（可多选）", 
                filetypes=[("图片", "*.png *.jpg *.jpeg *.bmp *.gif *.tiff"), ("所有", "*.*")]
            )
            if files:
                self.webp_files = list(files)
                self.webp_input.set(f"已选择 {len(files)} 个文件")
                self.webp_count_label.config(text=f"已选择 {len(files)} 张图片")
        else:
            folder = filedialog.askdirectory(title="选择文件夹")
            if folder:
                self.webp_input.set(folder)
                count = sum(1 for f in os.listdir(folder) 
                           if Path(f).suffix.lower() in {'.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff'})
                self.webp_count_label.config(text=f"文件夹中有 {count} 张图片")
    
    # ================ 转换方法 ================
    
    def convert_images_to_pdf(self):
        if not self.selected_files:
            messagebox.showwarning("提示", "请先添加图片")
            return
        output = filedialog.asksaveasfilename(
            title="保存PDF", defaultextension=".pdf", filetypes=[("PDF", "*.pdf")]
        )
        if output:
            self.converter.images_to_pdf(self.selected_files, output, self.image_quality.get())
    
    def convert_images_to_ppt(self):
        if not self.selected_files:
            messagebox.showwarning("提示", "请先添加图片")
            return
        output = filedialog.asksaveasfilename(
            title="保存PPT", defaultextension=".pptx", filetypes=[("PPT", "*.pptx")]
        )
        if output:
            self.converter.images_to_ppt(self.selected_files, output, self.image_quality.get())
    
    def convert_docs_to_pdf(self):
        if not self.doc_files:
            messagebox.showwarning("提示", "请先添加文档")
            return
        output = filedialog.askdirectory(title="选择输出文件夹")
        if output:
            self.converter.documents_to_pdf(self.doc_files, output)
    
    def convert_sheets_to_pdf(self):
        if not self.sheet_files:
            messagebox.showwarning("提示", "请先添加表格")
            return
        output = filedialog.askdirectory(title="选择输出文件夹")
        if output:
            self.converter.spreadsheets_to_pdf(self.sheet_files, output)
    
    def convert_urls_to_pdf(self):
        text = self.url_text.get(1.0, tk.END).strip()
        if not text:
            messagebox.showwarning("提示", "请输入网页地址")
            return
        
        urls = []
        for line in text.split('\n'):
            url = line.strip()
            if url and url != "https://":
                if not url.startswith(('http://', 'https://')):
                    url = 'https://' + url
                urls.append(url)
        
        if not urls:
            messagebox.showwarning("提示", "请输入有效的网页地址")
            return
        
        if len(urls) == 1:
            output = filedialog.asksaveasfilename(
                title="保存PDF", defaultextension=".pdf", filetypes=[("PDF", "*.pdf")]
            )
            if output:
                self.converter.url_to_pdf(urls[0], output)
        else:
            output = filedialog.askdirectory(title="选择输出文件夹")
            if output:
                self.converter.urls_to_pdf(urls, output)
    
    def convert_pdfs_to_ppt(self):
        if not self.pdf_files:
            messagebox.showwarning("提示", "请先添加PDF")
            return
        
        if len(self.pdf_files) == 1:
            output = filedialog.asksaveasfilename(
                title="保存PPT", defaultextension=".pptx", filetypes=[("PPT", "*.pptx")]
            )
            if output:
                self.converter.pdf_to_ppt(self.pdf_files[0], output, self.pdf_dpi.get())
        else:
            output = filedialog.askdirectory(title="选择输出文件夹")
            if output:
                self.converter.pdfs_to_ppt(self.pdf_files, output, self.pdf_dpi.get())
    
    def convert_pdfs_to_images(self):
        if not self.pdf_files:
            messagebox.showwarning("提示", "请先添加PDF")
            return
        output = filedialog.askdirectory(title="选择输出文件夹")
        if output:
            if len(self.pdf_files) == 1:
                self.converter.pdf_to_images(
                    self.pdf_files[0], output, self.pdf_dpi.get(), self.image_format.get()
                )
            else:
                self.converter.pdfs_to_images(
                    self.pdf_files, output, self.pdf_dpi.get(), self.image_format.get()
                )
    
    def extract_pdfs_images(self):
        if not self.pdf_files:
            messagebox.showwarning("提示", "请先添加PDF")
            return
        output = filedialog.askdirectory(title="选择输出文件夹")
        if output:
            if len(self.pdf_files) == 1:
                self.converter.extract_images_from_pdf(self.pdf_files[0], output)
            else:
                self.converter.extract_images_from_pdfs(self.pdf_files, output)
    
    def convert_to_webp(self):
        mode = self.webp_mode.get()
        if mode == 'files':
            if not self.webp_files:
                messagebox.showwarning("提示", "请先选择图片")
                return
            input_data = self.webp_files
        else:
            folder = self.webp_input.get()
            if not folder or not os.path.isdir(folder):
                messagebox.showwarning("提示", "请先选择文件夹")
                return
            input_data = folder
        
        output = filedialog.askdirectory(title="选择输出文件夹")
        if output:
            quality = self.webp_quality.get()
            resize = self.webp_resize.get()
            if mode == 'files':
                self.converter.images_to_webp(input_data, output, quality, resize)
            else:
                self.converter.folder_to_webp(input_data, output, quality, resize)


def main():
    import warnings
    warnings.filterwarnings('ignore')
    
    root = tk.Tk()
    app = ConverterGUI(root)
    root.mainloop()


if __name__ == "__main__":
    main()
